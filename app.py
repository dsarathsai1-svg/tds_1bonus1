# This is the complete Python code for your backend.
# Replace the entire content of `decksmith-ai/backend/app.py` with this code.

import os
import io
import json
from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from dotenv import load_dotenv
import google.generativeai as genai
from pptx import Presentation

# Load environment variables from .env file (for your GOOGLE_API_KEY)
load_dotenv()

# Configure the Flask app and enable Cross-Origin Resource Sharing (CORS)
app = Flask(__name__)
CORS(app)

# --- Helper Functions ---

def analyze_template(template_file):
    """
    Analyzes the uploaded .pptx template to extract slide layout names.
    Returns a list of layout names.
    """
    try:
        prs = Presentation(template_file)
        layout_names = [layout.name for layout in prs.slide_layouts]
        return layout_names
    except Exception as e:
        print(f"Error analyzing template: {e}")
        return None

def generate_presentation_structure(api_key, text, guidance, layout_names):
    """
    Calls the Gemini API to structure the text into a presentation format (JSON).
    """
    try:
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel('gemini-1.5-flash')

        # The master prompt that instructs the AI
        prompt = f"""
        You are an expert presentation designer. Your task is to convert the following text into a structured presentation format.

        Here is the list of available slide layouts you MUST use: {layout_names}

        User's Guidance for the tone/structure: "{guidance}"

        Raw Text to process:
        ---
        {text}
        ---

        Based on the text and guidance, create a JSON object for the presentation. The JSON should have a single key "slides", which is an array of slide objects.
        Each slide object must have two keys:
        1. "layout_name": A string that EXACTLY MATCHES one of the available layout names.
        2. "content": An object containing the text for the slide's placeholders. Common placeholders are 'title', 'subtitle', and 'points' (which should be an array of strings for bullet points). Choose placeholders that logically fit the chosen layout. For a 'Title Slide' layout, use 'title' and 'subtitle'. For a 'Title and Content' layout, use 'title' and 'points'.

        Example Output Format:
        {{
          "slides": [
            {{
              "layout_name": "Title Slide",
              "content": {{
                "title": "Presentation Title",
                "subtitle": "A Brief Subtitle"
              }}
            }},
            {{
              "layout_name": "Title and Content",
              "content": {{
                "title": "First Main Point",
                "points": [
                  "Detail point 1.",
                  "Detail point 2.",
                  "Detail point 3."
                ]
              }}
            }}
          ]
        }}

        Now, generate the JSON for the provided text. Output ONLY the raw JSON object, with no surrounding text or markdown formatting.
        """

        response = model.generate_content(prompt)
        # Clean up the response to ensure it's valid JSON
        cleaned_response = response.text.strip().replace('```json', '').replace('```', '')
        return json.loads(cleaned_response)
    except Exception as e:
        print(f"Error calling Gemini API: {e}")
        return None

def create_presentation(template_file, structure_json):
    """
    Creates a new .pptx file based on the template and the JSON structure.
    Returns the presentation file in a memory buffer.
    """
    try:
        prs = Presentation(template_file)
        slide_layouts = {layout.name: layout for layout in prs.slide_layouts}

        for slide_data in structure_json['slides']:
            layout_name = slide_data.get('layout_name')
            content = slide_data.get('content', {})

            if layout_name in slide_layouts:
                slide = prs.slides.add_slide(slide_layouts[layout_name])

                # Populate placeholders
                for shape in slide.placeholders:
                    placeholder_name = shape.name.lower()
                    if 'title' in placeholder_name and 'title' in content:
                        shape.text = content['title']
                    elif 'subtitle' in placeholder_name and 'subtitle' in content:
                        shape.text = content['subtitle']
                    elif 'body' in placeholder_name or 'content' in placeholder_name:
                        if 'points' in content and isinstance(content['points'], list):
                            tf = shape.text_frame
                            tf.clear()
                            p = tf.paragraphs[0]
                            p.text = content['points'][0]
                            for point_text in content['points'][1:]:
                                p = tf.add_paragraph()
                                p.text = point_text

        # Save the presentation to a memory buffer
        file_stream = io.BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)
        return file_stream
    except Exception as e:
        print(f"Error creating presentation: {e}")
        return None

# --- API Route ---

@app.route('/api/generate', methods=['POST'])
def generate_deck():
    """
    The main API endpoint to generate the presentation.
    """
    # 1. Get data from the request
    if 'template' not in request.files:
        return jsonify({"error": "No template file provided"}), 400

    template_file = request.files['template']
    text = request.form.get('text')
    guidance = request.form.get('guidance')
    api_key = request.form.get('apiKey')

    if not all([text, api_key]):
        return jsonify({"error": "Missing text or API key"}), 400

    # 2. Analyze the template
    layout_names = analyze_template(template_file)
    if not layout_names:
        return jsonify({"error": "Could not analyze the provided template file."}), 500
    
    # Rewind file pointer after analysis
    template_file.seek(0)

    # 3. Call LLM to get the presentation structure
    structure_json = generate_presentation_structure(api_key, text, guidance, layout_names)
    if not structure_json:
        return jsonify({"error": "Failed to generate presentation structure from the language model."}), 500

    # 4. Create the new presentation
    presentation_stream = create_presentation(template_file, structure_json)
    if not presentation_stream:
        return jsonify({"error": "Failed to create the final presentation file."}), 500

    # 5. Send the file back to the user for download
    return send_file(
        presentation_stream,
        as_attachment=True,
        download_name='generated_presentation.pptx',
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)

# --- How to Run and Test ---
# 1. Make sure your `.env` file is in `decksmith-ai/backend` and contains your GOOGLE_API_KEY.
# 2. In your terminal, inside `decksmith-ai/backend`, run: python app.py
# 3. The server is now running. We will test it from the frontend in the next steps.